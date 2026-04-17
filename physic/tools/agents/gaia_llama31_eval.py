import argparse
import datetime as _dt
import hashlib
import json
import os
import re
import sys
import time
import urllib.error
import urllib.request
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Tuple


def _read_jsonl(path: Path) -> List[Dict[str, Any]]:
    rows: List[Dict[str, Any]] = []
    with path.open("r", encoding="utf-8") as f:
        for line in f:
            line = line.strip()
            if not line:
                continue
            rows.append(json.loads(line))
    return rows


def _norm(s: str) -> str:
    s2 = (s or "").strip().lower()
    s2 = re.sub(r"\s+", " ", s2)
    return s2


def _norm_relaxed(s: str) -> str:
    s2 = _norm(s)
    s2 = re.sub(r"[^\w\s]", "", s2)
    s2 = re.sub(r"\s+", " ", s2).strip()
    return s2


def _clip(s: str, max_chars: int) -> str:
    if len(s) <= max_chars:
        return s
    head = s[: max_chars // 2]
    tail = s[-max_chars // 2 :]
    return head + "\n...(clipped)...\n" + tail


def _sha256_file(path: Path, max_bytes: int = 2_000_000) -> str:
    h = hashlib.sha256()
    n = 0
    with path.open("rb") as f:
        while True:
            chunk = f.read(65536)
            if not chunk:
                break
            h.update(chunk)
            n += len(chunk)
            if n >= max_bytes:
                break
    return h.hexdigest()


def _read_text_file(path: Path, max_chars: int) -> str:
    try:
        raw = path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        raw = path.read_text(encoding="utf-8", errors="replace")
    return _clip(raw, max_chars)


def _read_pdf(path: Path, max_chars: int) -> str:
    try:
        from pypdf import PdfReader  # type: ignore

        reader = PdfReader(str(path))
        parts: List[str] = []
        for p in reader.pages:
            try:
                parts.append(p.extract_text() or "")
            except Exception:
                parts.append("")
        return _clip("\n".join(parts).strip(), max_chars)
    except Exception:
        return ""


def _read_docx(path: Path, max_chars: int) -> str:
    try:
        import docx  # type: ignore

        d = docx.Document(str(path))
        parts = [p.text for p in d.paragraphs if (p.text or "").strip()]
        return _clip("\n".join(parts).strip(), max_chars)
    except Exception:
        return ""


def _read_pptx(path: Path, max_chars: int) -> str:
    try:
        from pptx import Presentation  # type: ignore

        prs = Presentation(str(path))
        parts: List[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                text = getattr(shape, "text", None)
                if text:
                    t = str(text).strip()
                    if t:
                        parts.append(t)
        return _clip("\n".join(parts).strip(), max_chars)
    except Exception:
        return ""


def _read_csv(path: Path, max_chars: int) -> str:
    try:
        import pandas as pd  # type: ignore

        df = pd.read_csv(path)
        return _clip(df.to_string(index=False), max_chars)
    except Exception:
        return _read_text_file(path, max_chars=max_chars)


def _read_xlsx(path: Path, max_chars: int) -> str:
    try:
        import pandas as pd  # type: ignore

        xls = pd.ExcelFile(path)
        parts: List[str] = []
        for name in xls.sheet_names[:5]:
            df = xls.parse(name)
            parts.append(f"### Sheet: {name}\n{df.to_string(index=False)}")
        return _clip("\n\n".join(parts).strip(), max_chars)
    except Exception:
        return ""


def _extract_attachment_text(path: Path, max_chars: int) -> Tuple[str, str]:
    suffix = path.suffix.lower()
    if suffix in {".txt", ".json", ".jsonl", ".xml", ".md", ".py"}:
        return _read_text_file(path, max_chars=max_chars), "text"
    if suffix in {".csv"}:
        return _read_csv(path, max_chars=max_chars), "csv"
    if suffix in {".xlsx"}:
        return _read_xlsx(path, max_chars=max_chars), "xlsx"
    if suffix in {".pdf"}:
        return _read_pdf(path, max_chars=max_chars), "pdf"
    if suffix in {".docx"}:
        return _read_docx(path, max_chars=max_chars), "docx"
    if suffix in {".pptx"}:
        return _read_pptx(path, max_chars=max_chars), "pptx"
    return "", "unknown"


def _ollama_generate(
    *,
    host: str,
    model: str,
    prompt: str,
    temperature: float,
    num_predict: int,
    timeout_s: int,
) -> str:
    url = host.rstrip("/") + "/api/generate"
    payload = {
        "model": model,
        "prompt": prompt,
        "stream": False,
        "options": {"temperature": float(temperature), "num_predict": int(num_predict)},
    }
    data = json.dumps(payload).encode("utf-8")
    req = urllib.request.Request(url, data=data, headers={"Content-Type": "application/json"})
    with urllib.request.urlopen(req, timeout=timeout_s) as resp:
        body = resp.read().decode("utf-8", errors="replace")
    obj = json.loads(body)
    return str(obj.get("response") or "").strip()


def _extract_final_answer(text: str) -> str:
    t = (text or "").strip()
    if not t:
        return ""
    m = re.search(r"\{[\s\S]*\}", t)
    if m:
        blob = m.group(0)
        try:
            obj = json.loads(blob)
            if isinstance(obj, dict):
                fa = obj.get("final_answer")
                if isinstance(fa, str):
                    return fa.strip()
        except Exception:
            pass
    t = re.sub(r"(?im)^\s*final\s*answer\s*:\s*", "", t).strip()
    t = t.splitlines()[0].strip()
    return t


def _build_prompt(question: str, attachment_block: str) -> str:
    return (
        "You are answering a GAIA benchmark question.\n"
        "Return ONLY valid JSON, exactly in this format:\n"
        '{"final_answer": "<answer>"}\n'
        "No extra keys, no markdown, no explanation.\n\n"
        f"Question:\n{question.strip()}\n\n"
        f"{attachment_block}"
    )


def _attachment_block(
    *,
    split_dir: Path,
    file_name: str,
    max_chars: int,
) -> str:
    fn = (file_name or "").strip()
    if not fn:
        return "Attachment: (none)\n"
    p = (split_dir / fn).resolve()
    if not p.exists():
        return f"Attachment: {fn} (missing)\n"
    extracted, kind = _extract_attachment_text(p, max_chars=max_chars)
    meta = {
        "file_name": fn,
        "kind": kind,
        "size_bytes": p.stat().st_size,
        "sha256_prefix": _sha256_file(p)[:12],
    }
    if extracted.strip():
        return (
            "Attachment metadata:\n"
            + json.dumps(meta, ensure_ascii=False)
            + "\n\nAttachment content (text extraction):\n"
            + extracted.strip()
            + "\n"
        )
    return (
        "Attachment metadata:\n"
        + json.dumps(meta, ensure_ascii=False)
        + "\n\nAttachment content: (not extractable by this script)\n"
    )


def _select_rows(
    rows: Iterable[Dict[str, Any]],
    *,
    levels: Optional[List[int]],
) -> List[Dict[str, Any]]:
    out: List[Dict[str, Any]] = []
    for r in rows:
        lv = r.get("Level")
        try:
            lv_int = int(lv)
        except Exception:
            continue
        if levels and (lv_int not in levels):
            continue
        out.append(r)
    return out


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--data_root", type=str, required=True)
    ap.add_argument("--year", type=str, default="2023")
    ap.add_argument("--split", type=str, choices=["validation", "test"], default="validation")
    ap.add_argument("--levels", type=str, default="1,2,3")
    ap.add_argument("--limit", type=int, default=0)
    ap.add_argument("--offset", type=int, default=0)
    ap.add_argument("--shuffle", action="store_true")
    ap.add_argument("--seed", type=int, default=0)

    ap.add_argument("--model", type=str, default="llama3.1:latest")
    ap.add_argument("--ollama_host", type=str, default=os.environ.get("OLLAMA_HOST", "http://127.0.0.1:11434"))
    ap.add_argument("--temperature", type=float, default=0.0)
    ap.add_argument("--num_predict", type=int, default=256)
    ap.add_argument("--timeout_s", type=int, default=180)
    ap.add_argument("--max_attachment_chars", type=int, default=12000)
    ap.add_argument("--max_question_chars", type=int, default=6000)
    ap.add_argument("--sleep_s", type=float, default=0.0)

    ap.add_argument("--out_dir", type=str, default="")
    ap.add_argument("--resume", action="store_true")
    ap.add_argument("--dry_run", action="store_true")
    args = ap.parse_args()

    data_root = Path(args.data_root).expanduser().resolve()
    split_dir = data_root / args.year / args.split
    meta_path = split_dir / "metadata.jsonl"
    if not meta_path.exists():
        raise FileNotFoundError(f"metadata.jsonl not found: {meta_path}")

    levels = [int(x) for x in re.split(r"[,\s]+", args.levels.strip()) if x.strip()]
    rows = _select_rows(_read_jsonl(meta_path), levels=levels)
    if args.shuffle:
        import random

        random.Random(args.seed).shuffle(rows)
    if args.offset:
        rows = rows[args.offset :]
    if args.limit and args.limit > 0:
        rows = rows[: args.limit]

    ts = _dt.datetime.now().strftime("%Y%m%d_%H%M%S")
    out_dir = Path(args.out_dir).expanduser().resolve() if args.out_dir else (Path(__file__).resolve().parent / "out_gaia_llama31")
    out_dir.mkdir(parents=True, exist_ok=True)
    pred_path = out_dir / f"predictions_{args.year}_{args.split}_{ts}.jsonl"
    summary_path = out_dir / f"summary_{args.year}_{args.split}_{ts}.json"

    done_task_ids: set[str] = set()
    if args.resume:
        existing = sorted(out_dir.glob(f"predictions_{args.year}_{args.split}_*.jsonl"))
        if existing:
            latest = existing[-1]
            with latest.open("r", encoding="utf-8") as f:
                for line in f:
                    try:
                        obj = json.loads(line)
                        tid = str(obj.get("task_id") or "").strip()
                        if tid:
                            done_task_ids.add(tid)
                    except Exception:
                        continue

    n_total = 0
    n_scored = 0
    n_em = 0
    n_relaxed = 0
    
    # Track stats by level
    level_stats = {
        1: {"n_scored": 0, "n_em": 0, "n_relaxed": 0},
        2: {"n_scored": 0, "n_em": 0, "n_relaxed": 0},
        3: {"n_scored": 0, "n_em": 0, "n_relaxed": 0}
    }

    with pred_path.open("w", encoding="utf-8") as out_f:
        for r in rows:
            task_id = str(r.get("task_id") or "").strip()
            if not task_id:
                continue
            if task_id in done_task_ids:
                continue
            question = str(r.get("Question") or "")
            answer = str(r.get("Final answer") or "").strip()
            level = r.get("Level")
            try:
                level_int = int(level)
            except (ValueError, TypeError):
                level_int = 1 # Fallback, though GAIA always has 1, 2, or 3
            file_name = str(r.get("file_name") or "")

            question = _clip(question, args.max_question_chars)
            att_block = _attachment_block(split_dir=split_dir, file_name=file_name, max_chars=args.max_attachment_chars)
            prompt = _build_prompt(question, att_block)

            if args.dry_run:
                model_text = ""
                final_answer = ""
                err = ""
            else:
                try:
                    model_text = _ollama_generate(
                        host=args.ollama_host,
                        model=args.model,
                        prompt=prompt,
                        temperature=args.temperature,
                        num_predict=args.num_predict,
                        timeout_s=args.timeout_s,
                    )
                    final_answer = _extract_final_answer(model_text)
                    err = ""
                except urllib.error.URLError as e:
                    model_text = ""
                    final_answer = ""
                    err = f"ollama_error: {e}"
                except Exception as e:
                    model_text = ""
                    final_answer = ""
                    err = f"error: {e}"

            em = False
            relaxed = False
            scored = bool(answer) and (answer != "?") and not args.dry_run
            if scored:
                em = _norm(final_answer) == _norm(answer)
                relaxed = _norm_relaxed(final_answer) == _norm_relaxed(answer)
                n_scored += 1
                n_em += int(em)
                n_relaxed += int(relaxed)
                if level_int in level_stats:
                    level_stats[level_int]["n_scored"] += 1
                    level_stats[level_int]["n_em"] += int(em)
                    level_stats[level_int]["n_relaxed"] += int(relaxed)

            rec = {
                "task_id": task_id,
                "level": level,
                "file_name": file_name,
                "gold": answer,
                "pred": final_answer,
                "raw": model_text,
                "em": em if scored else None,
                "relaxed_em": relaxed if scored else None,
                "error": err or None,
            }
            out_f.write(json.dumps(rec, ensure_ascii=False) + "\n")
            out_f.flush()

            n_total += 1
            if args.sleep_s and (not args.dry_run):
                time.sleep(float(args.sleep_s))

    summary = {
        "data_root": str(data_root),
        "year": args.year,
        "split": args.split,
        "levels": levels,
        "model": args.model,
        "ollama_host": args.ollama_host,
        "predictions_path": str(pred_path),
        "n_total": n_total,
        "n_scored": n_scored,
        "em": (n_em / n_scored) if n_scored else None,
        "relaxed_em": (n_relaxed / n_scored) if n_scored else None,
        "level_stats": {
            k: {
                "n_scored": v["n_scored"],
                "em": (v["n_em"] / v["n_scored"]) if v["n_scored"] else None,
                "relaxed_em": (v["n_relaxed"] / v["n_scored"]) if v["n_scored"] else None,
            } for k, v in level_stats.items() if v["n_scored"] > 0
        },
        "timestamp": ts,
    }
    summary_path.write_text(json.dumps(summary, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

    print(json.dumps(summary, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

